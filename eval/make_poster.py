"""
Research poster generator — Atlas (DATA 298B, Team 1)
Output: eval/poster/atlas_poster.pptx   (48" × 36")

Run: python eval/make_poster.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours (SJSU palette) ──────────────────────────────────────────────
NAVY    = RGBColor(0x00, 0x22, 0x55)   # #002255 — header / accents
ORANGE  = RGBColor(0xE5, 0x72, 0x00)   # #E57200 — section bars / strips
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF2, 0xF4, 0xF6)   # panel backgrounds
DKGRAY  = RGBColor(0x33, 0x33, 0x33)   # body text
MIDGRAY = RGBColor(0x66, 0x66, 0x66)   # secondary text
GOLD    = RGBColor(0xFF, 0xC0, 0x00)   # highlight in header

# ── Poster dimensions ─────────────────────────────────────────────────
W_IN = 48.0   # width  (inches)
H_IN = 36.0   # height (inches)

# gutter (decorative orange strip width on each side)
GUTTER_IN = 0.9

# header height
HDR_H_IN = 5.2

# footer height
FTR_H_IN = 0.7

# column layout  (3 columns inside the gutters)
COL_X0_IN = GUTTER_IN + 0.3       # first column left edge
COL_AREA_W_IN = W_IN - 2 * (GUTTER_IN + 0.3) - 0.1
COL_GAP_IN = 0.35
COL_W_IN = (COL_AREA_W_IN - 2 * COL_GAP_IN) / 3

CONTENT_Y0_IN = HDR_H_IN + 0.35   # top of first row of content
CONTENT_H_IN  = H_IN - HDR_H_IN - FTR_H_IN - 0.55

# ── Helpers ───────────────────────────────────────────────────────────
def col_x(n: int) -> float:
    """Left-edge X of column n (0-indexed) in inches."""
    return COL_X0_IN + n * (COL_W_IN + COL_GAP_IN)


def add_rect(slide, x, y, w, h, fill: RGBColor, border=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=10, bold=False,
                color=DKGRAY, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def section_box(slide, x, y, w, h, title, body_lines,
                title_size=11.5, body_size=9.5):
    """Orange title bar + white content panel."""
    bar_h = 0.32
    add_rect(slide, x, y, w, bar_h, ORANGE)
    add_textbox(slide, x + 0.07, y + 0.01, w - 0.14, bar_h,
                title, size=title_size, bold=True, color=WHITE)
    panel_h = h - bar_h
    add_rect(slide, x, y + bar_h, w, panel_h, LGRAY)
    body_text = "\n".join(body_lines)
    add_textbox(slide, x + 0.12, y + bar_h + 0.1,
                w - 0.24, panel_h - 0.15,
                body_text, size=body_size, color=DKGRAY)
    return y + h


def bullet_section(slide, x, y, w, h, title, bullets,
                   title_size=11.5, body_size=9.2):
    """Section box where bullets are formatted with • prefix."""
    lines = [f"•  {b}" for b in bullets]
    return section_box(slide, x, y, w, h, title, lines,
                       title_size=title_size, body_size=body_size)


def add_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


# ── Chart paths ──────────────────────────────────────────────────────
CHARTS = {
    "radar":     "eval/charts/1_radar_dimensions.png",
    "cat_bar":   "eval/charts/2_geval_by_category.png",
    "dual":      "eval/charts/3_dual_metric_by_category.png",
    "scatter":   "eval/charts/4_scatter_task_vs_geval.png",
    "heatmap":   "eval/charts/5_heatmap_dim_category.png",
    "dist":      "eval/charts/6_geval_distribution.png",
    "donut":     "eval/charts/7_category_distribution.png",
    "summary":   "eval/charts/8_summary_dashboard.png",
}

# ── Build poster ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(W_IN)
prs.slide_height = Inches(H_IN)

slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

# ── Background (white) ───────────────────────────────────────────────
add_rect(slide, 0, 0, W_IN, H_IN, WHITE)

# ── Left & right orange gutter strips ───────────────────────────────
add_rect(slide, 0,           0, GUTTER_IN, H_IN, ORANGE)
add_rect(slide, W_IN - GUTTER_IN, 0, GUTTER_IN, H_IN, ORANGE)

# ── Header ───────────────────────────────────────────────────────────
add_rect(slide, GUTTER_IN, 0, W_IN - 2 * GUTTER_IN, HDR_H_IN, NAVY)

# University + dept
add_textbox(slide,
            GUTTER_IN + 0.3, 0.18,
            W_IN - 2 * GUTTER_IN - 0.6, 0.6,
            "San José State University  ·  College of Science  ·  Data Science M.S. Program (DATA 298B)",
            size=13, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

# Project title
add_textbox(slide,
            GUTTER_IN + 0.3, 0.75,
            W_IN - 2 * GUTTER_IN - 0.6, 1.6,
            "Atlas: AI-Powered Customer Support Agent\n"
            "with Fine-Tuned Phi-4-mini and RAG Pipeline",
            size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Team + advisor
add_textbox(slide,
            GUTTER_IN + 0.3, 2.50,
            W_IN - 2 * GUTTER_IN - 0.6, 0.55,
            "Thota Himaja Sree   ·   Rondla Chandana   ·   Patel Meetkumar   ·   Cheedalla Suharsha",
            size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide,
            GUTTER_IN + 0.3, 3.05,
            W_IN - 2 * GUTTER_IN - 0.6, 0.45,
            "Faculty Advisor: Prof. Simon Shim     ·     Team 1",
            size=14, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

# Thin gold rule under header
add_rect(slide, GUTTER_IN, HDR_H_IN - 0.12, W_IN - 2 * GUTTER_IN, 0.12, ORANGE)

# Thin gold rule at very top inside strip (decoration)
add_rect(slide, 0, HDR_H_IN - 0.12, W_IN, 0.12, ORANGE)

# ── Column X coordinates ─────────────────────────────────────────────
cx = [col_x(i) for i in range(3)]

# ── ROW HEIGHTS — three rows inside CONTENT_H_IN ────────────────────
# Row 1: Introduction + Architecture + Methodology
# Row 2: Fine-Tuning + Evaluation design + Key Results charts
# Row 3: G-Eval breakdown charts + Conclusions + References

R1 = 9.6
R2 = 10.2
R3 = CONTENT_H_IN - R1 - R2 - 0.4   # ~9.2"

y1 = CONTENT_Y0_IN
y2 = y1 + R1 + 0.18
y3 = y2 + R2 + 0.22

# ═══════════════════════════════════════════════════════════
# COLUMN 1
# ═══════════════════════════════════════════════════════════

# --- Introduction ---
bullet_section(slide, cx[0], y1, COL_W_IN, 4.3,
    "1. Introduction",
    [
        "E-commerce platforms require scalable, intelligent customer support",
        "Rule-based chatbots fail on long-tail queries and nuanced language",
        "LLMs offer fluent responses but hallucinate without grounding",
        "Atlas combines fine-tuned SLM + RAG + deterministic tool routing",
        "Goal: accurate, empathetic, grounded support — at small-model cost",
    ])

# --- Motivation ---
bullet_section(slide, cx[0], y1 + 4.45, COL_W_IN, 2.55,
    "2. Motivation & Problem Statement",
    [
        "Support tickets cost ~$5–15/contact; automation saves 40–70%",
        "Customers demand instant resolution of order, refund, and policy queries",
        "Challenge: hallucination, privacy (ID verification), latency, cost",
        "Research question: Can a fine-tuned 3.8B SLM + RAG match GPT-4-class accuracy?",
    ])

# --- Dataset ---
bullet_section(slide, cx[0], y1 + 7.15, COL_W_IN, 2.28,
    "3. Dataset",
    [
        "Amazon product catalog — 359K items (title, category, price, description)",
        "Policy knowledge base — 45 hand-authored entries (returns, refunds, shipping)",
        "Fine-tune corpus — 2,400 synthetic customer-support Q&A pairs",
        "Evaluation set — 50 hand-crafted test cases across 8 intent categories",
    ])

# Row 2 Col 1 — System Architecture
bullet_section(slide, cx[0], y2, COL_W_IN, 4.8,
    "4. System Architecture",
    [
        "FastAPI backend · PostgreSQL (orders, sessions) · Qdrant vector DB",
        "HuggingFace Inference Endpoint hosting fine-tuned Phi-4-mini adapter",
        "Retrieval: sentence-transformers/all-MiniLM-L6-v2 (384-dim) · top-6 cosine",
        "Deterministic tool router (keyword/regex) — 8 tools:",
        "    lookup_order · process_refund · cancel_order · get_account_info",
        "    search_product/policy_knowledge · escalate_to_human · send_email",
        "Identity verification gate — email/order-ID required before PII access",
        "Two-layer LLM safety net: stall-detection + data-presence → template fallback",
        "Frontend: vanilla JS chat UI · MailHog / Gmail SMTP email delivery",
        "Deployment: Docker Compose  (api · postgres · qdrant · mailhog)",
    ])

# Row 2 Col 1 continued — Fine-Tuning
bullet_section(slide, cx[0], y2 + 5.0, COL_W_IN, 5.05,
    "5. Model Selection & Fine-Tuning",
    [
        "Evaluated 4 open-weight models on 20 qualitative prompts:",
        "  LLaMA-3.2-1B · SmolLM2-1.7B · Qwen2.5-1.5B · Phi-4-mini-instruct",
        "Selected: Phi-4-mini-instruct (3.8B) — best instruction following & tone",
        "Fine-tune method: QLoRA (4-bit NF4 quantization) on Google Colab A100",
        "Training data: 2,400 support Q&A pairs · 3 epochs · batch 4 · lr 2e-4",
        "Adapter uploaded to HuggingFace Hub; served on A10G GPU endpoint",
        "Inference: vLLM-compatible · avg first-token latency ~3–5 s (warm)",
        "Cold-start after idle auto-pause: 30–60 s warm-up",
    ])

# ═══════════════════════════════════════════════════════════
# COLUMN 2
# ═══════════════════════════════════════════════════════════

# --- Evaluation Design ---
bullet_section(slide, cx[1], y1, COL_W_IN, 4.9,
    "6. Evaluation Framework",
    [
        "50 hand-crafted test cases spanning 8 intent categories:",
        "  order lookup · refund · cancellation · product recommendation",
        "  policy FAQ · escalation · account info · general greeting",
        "Metrics:",
        "  • ROUGE-L — surface lexical overlap with reference answers",
        "  • Task Success Rate — binary: correct tool call + grounded response",
        "  • G-Eval (1–5) — Gemini 2.5 Flash as LLM judge across 5 dimensions:",
        "      Relevance · Faithfulness · Completeness · Tone & Empathy · Groundedness",
        "Identity-sensitive cases pre-seeded customer_id to bypass ID gate",
        "Evaluation pipeline: evaluate.py · results stored in results.md",
    ])

# Overall metrics table
section_box(slide, cx[1], y1 + 5.05, COL_W_IN, 4.42,
    "7. Overall Evaluation Results",
    [
        "  Metric                        Score",
        "  ─────────────────────────────────────",
        "  ROUGE-L                       0.191",
        "  Task Success Rate             72.3 %",
        "  G-Eval Average (1–5)          3.79",
        "  ─────────────────────────────────────",
        "  Relevance                     4.24",
        "  Tone & Empathy                4.18",
        "  Faithfulness                  3.78",
        "  Groundedness                  3.44",
        "  Completeness                  3.30",
        "  ─────────────────────────────────────",
        "  Best category: Escalation     4.80",
        "  Weakest category: Cancellation 2.73",
    ],
    body_size=9.4)

# Row 2 Col 2 — Radar chart
bar_h2 = 0.32
add_rect(slide, cx[1], y2, COL_W_IN, bar_h2, ORANGE)
add_textbox(slide, cx[1] + 0.07, y2 + 0.01, COL_W_IN - 0.14, bar_h2,
            "8. G-Eval Dimension Radar", size=11.5, bold=True, color=WHITE)
add_image(slide, CHARTS["radar"], cx[1], y2 + bar_h2,
          COL_W_IN, 4.85)

# Row 2 Col 2 — Heatmap
bar_h3 = 0.32
hm_y = y2 + bar_h2 + 4.85 + 0.12
add_rect(slide, cx[1], hm_y, COL_W_IN, bar_h3, ORANGE)
add_textbox(slide, cx[1] + 0.07, hm_y + 0.01, COL_W_IN - 0.14, bar_h3,
            "9. G-Eval Heatmap: Dimension × Category", size=11.5, bold=True, color=WHITE)
add_image(slide, CHARTS["heatmap"], cx[1], hm_y + bar_h3,
          COL_W_IN, R2 - bar_h2 - 4.85 - 0.12 - bar_h3 - 0.05)

# Row 3 Col 2 — G-Eval by category bar
bar_h4 = 0.32
add_rect(slide, cx[1], y3, COL_W_IN, bar_h4, ORANGE)
add_textbox(slide, cx[1] + 0.07, y3 + 0.01, COL_W_IN - 0.14, bar_h4,
            "10. G-Eval by Intent Category", size=11.5, bold=True, color=WHITE)
img_h_cat = min(R3 * 0.54, 4.6)
add_image(slide, CHARTS["cat_bar"], cx[1], y3 + bar_h4,
          COL_W_IN, img_h_cat)

# Conclusions
conc_y = y3 + bar_h4 + img_h_cat + 0.12
conc_h = R3 - bar_h4 - img_h_cat - 0.12
bullet_section(slide, cx[1], conc_y, COL_W_IN, conc_h,
    "13. Conclusions",
    [
        "Fine-tuned Phi-4-mini + RAG achieves 72.3% task success on e-commerce support",
        "Deterministic routing + template fallback critical for reliability (ROUGE-L limited)",
        "Escalation best handled (4.80); cancellation most problematic (2.73)",
        "Identity verification gate prevents data leakage with minimal UX friction",
        "RAG grounding reduces hallucination vs. base LLM; re-ranking would further help",
    ],
    body_size=9.2)

# ═══════════════════════════════════════════════════════════
# COLUMN 3
# ═══════════════════════════════════════════════════════════

# Row 1 Col 3 — Dual metric by category
bar_h5 = 0.32
add_rect(slide, cx[2], y1, COL_W_IN, bar_h5, ORANGE)
add_textbox(slide, cx[2] + 0.07, y1 + 0.01, COL_W_IN - 0.14, bar_h5,
            "11. Task Success & G-Eval by Category", size=11.5, bold=True, color=WHITE)
add_image(slide, CHARTS["dual"], cx[2], y1 + bar_h5, COL_W_IN, 4.65)

# Task Success scatter
bar_h6 = 0.32
sc_y = y1 + bar_h5 + 4.65 + 0.12
add_rect(slide, cx[2], sc_y, COL_W_IN, bar_h6, ORANGE)
add_textbox(slide, cx[2] + 0.07, sc_y + 0.01, COL_W_IN - 0.14, bar_h6,
            "12. Task Success vs G-Eval Score (per case)", size=11.5, bold=True, color=WHITE)
add_image(slide, CHARTS["scatter"], cx[2], sc_y + bar_h6,
          COL_W_IN, R1 - bar_h5 - 4.65 - 0.12 - bar_h6 - 0.04)

# Row 2 Col 3 — Summary dashboard
bar_h7 = 0.32
add_rect(slide, cx[2], y2, COL_W_IN, bar_h7, ORANGE)
add_textbox(slide, cx[2] + 0.07, y2 + 0.01, COL_W_IN - 0.14, bar_h7,
            "Summary Dashboard", size=11.5, bold=True, color=WHITE)
add_image(slide, CHARTS["summary"], cx[2], y2 + bar_h7,
          COL_W_IN, R2 - bar_h7 - 0.05)

# Row 3 Col 3 — Future work + References + Ack
bullet_section(slide, cx[2], y3, COL_W_IN, 4.25,
    "14. Future Work & Limitations",
    [
        "Deploy no-RAG baseline (Variant B) for direct ablation comparison",
        "Add cross-encoder re-ranking over Qdrant top-k results",
        "Streaming responses (Server-Sent Events) to reduce perceived latency",
        "Stronger identity verification: OTP or OAuth2 instead of email-only",
        "Multi-turn context compression for conversations > 10 turns",
        "Cancellation flow refinement — communicate post-delivery return path",
        "Mobile-responsive frontend and persistent session storage",
    ],
    body_size=9.2)

# References
section_box(slide, cx[2], y3 + 4.38, COL_W_IN, 2.55,
    "15. References",
    [
        "[1] Dettmers et al. (2023). QLoRA: Efficient Finetuning of Quantized LLMs. NeurIPS.",
        "[2] Lewis et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP. NeurIPS.",
        "[3] Liu et al. (2023). G-Eval: NLG Evaluation using GPT-4 with Better Human Alignment. EMNLP.",
        "[4] Microsoft Phi-4-mini-instruct. HuggingFace Hub, 2024.",
        "[5] Amazon Product Reviews Dataset (2023). Kaggle.",
        "[6] Sentence-Transformers: all-MiniLM-L6-v2. HuggingFace, 2021.",
    ],
    body_size=8.5)

# Acknowledgements
section_box(slide, cx[2], y3 + 7.06, COL_W_IN, CONTENT_H_IN - (y3 - CONTENT_Y0_IN) - 7.06 + CONTENT_Y0_IN - y3 + (R3 - 7.06),
    "Acknowledgements",
    [
        "We thank Prof. Simon Shim for guidance throughout DATA 298B.",
        "Google Colab provided A100 GPU access for fine-tuning experiments.",
        "HuggingFace Inference Endpoints (A10G) hosted the production model.",
        "GitHub: github.com/Ch-Suharsha/teammate-rag",
    ],
    body_size=9.0)

# ── Footer ────────────────────────────────────────────────────────────
ftr_y = H_IN - FTR_H_IN
add_rect(slide, GUTTER_IN, ftr_y, W_IN - 2 * GUTTER_IN, FTR_H_IN, NAVY)
add_textbox(slide,
            GUTTER_IN + 0.3, ftr_y + 0.1,
            W_IN - 2 * GUTTER_IN - 0.6, FTR_H_IN - 0.1,
            "DATA 298B Master's Project  ·  San José State University  ·  Spring 2025  ·  "
            "Advisor: Prof. Simon Shim  ·  github.com/Ch-Suharsha/teammate-rag",
            size=11, color=GOLD, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────
out_dir = "eval/poster"
os.makedirs(out_dir, exist_ok=True)
pptx_path = f"{out_dir}/atlas_poster.pptx"
prs.save(pptx_path)
print(f"Saved → {pptx_path}")
print("Open in PowerPoint (or Keynote / LibreOffice Impress) and export to PDF.")
